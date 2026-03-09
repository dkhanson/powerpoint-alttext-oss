import io
from pptx import Presentation
from pptx.util import Inches

from powerpoint_alttext_v2.core.processor import PowerPointProcessor


def create_test_presentation():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Add three textboxes in different positions
    tb1 = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(1))
    tb1.text_frame.text = "First"

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(2), Inches(1))
    tb2.text_frame.text = "Second"

    tb3 = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(2), Inches(1))
    tb3.text_frame.text = "Third"

    return prs


def test_annotate_reading_order_sets_title_tags(tmp_path):
    prs = create_test_presentation()
    slide = prs.slides[0]

    processor = PowerPointProcessor()
    # Ensure no pre-existing title attributes
    for shp in slide.shapes:
        c = getattr(shp, '_element', None)
        # Remove any title attr for a clean test
        try:
            nv = c.find('{http://schemas.openxmlformats.org/presentationml/2006/main}nvSpPr')
            if nv is not None:
                cNvPr = nv.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
                if cNvPr is not None and cNvPr.get('title'):
                    cNvPr.set('title', '')
        except Exception:
            pass

    processor.annotate_reading_order(slide)

    # After annotation, the shapes should have title attributes containing 'ord 1', 'ord 2', 'ord 3'
    titles = []
    for shp in slide.shapes:
        c = getattr(shp, '_element', None)
        nv = c.find('{http://schemas.openxmlformats.org/presentationml/2006/main}nvSpPr')
        cNvPr = nv.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
        titles.append((shp, cNvPr.get('title') or ''))

    # Titles should contain ord 1..3 in some order
    found = [t for _, t in titles]
    assert any('ord 1' in s for s in found)
    assert any('ord 2' in s for s in found)
    assert any('ord 3' in s for s in found)
