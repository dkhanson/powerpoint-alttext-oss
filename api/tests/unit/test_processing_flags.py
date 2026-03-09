from pptx import Presentation
from pptx.util import Inches

from powerpoint_alttext_v2.core.processor import PowerPointProcessor
from powerpoint_alttext_v2.config import ProcessingConfig


def create_slide_with_textbox():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(1))
    tb.text_frame.text = "Hello"
    return prs, slide, tb


def test_images_and_decorators_only_skips_textboxes(monkeypatch):
    prs, slide, tb = create_slide_with_textbox()
    processor = PowerPointProcessor()

    # Patch config to enable the flag
    processor.config.processing.images_and_decorators_only = True

    # Process the textbox shape
    result = processor.process_shape(tb, slide_num=1, shape_idx=0)

    # Expect it to be skipped (not processed)
    assert result['processed'] is False


def test_skip_objects_with_text_skips_shapes_with_text(monkeypatch):
    prs, slide, tb = create_slide_with_textbox()
    processor = PowerPointProcessor()

    # Patch config to enable skip_objects_with_text
    processor.config.processing.skip_objects_with_text = True

    result = processor.process_shape(tb, slide_num=1, shape_idx=0)
    assert result['processed'] is False
