"""
PowerPoint Alt-Text Generator Core Processor (V2)
Based on the working main_no_com.py script
"""

import os
import sys
import base64
import io
import shutil
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import List, Tuple, Optional, Dict, Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn, _nsmap
from pptx.oxml.xmlchemy import OxmlElement

# Image processing for memory optimization
try:
    from PIL import Image
except ImportError:
    Image = None

try:
    from openai import AzureOpenAI, OpenAI
except Exception:
    AzureOpenAI = None
    OpenAI = None

from ..config import get_config

# Namespaces and extension IDs for decorative flag
ADEC_NS = "http://schemas.microsoft.com/office/drawing/2017/decorative"
DECORATIVE_EXT_URI = "{C183D7F6-B498-43B3-948B-1728B52AA6E4}"

# Ensure adec prefix is known for XPath / element creation
if "adec" not in _nsmap:
    _nsmap["adec"] = ADEC_NS


class PowerPointProcessor:
    """PowerPoint Alt-Text Generator using TOML configuration."""
    
    def __init__(self):
        self.config = get_config()
        self.client = None
        self._seen_titles = {}
        self._image_cache = {}  # Cache for duplicate images
        
        # Initialize OpenAI client (Azure or standard)
        if self.config.azure_openai.api_key:
            try:
                timeout = getattr(self.config.processing, 'api_timeout_seconds', 30)
                max_retries = getattr(self.config.processing, 'max_retries', 2)
                provider = getattr(self.config.azure_openai, 'provider', 'auto')

                use_azure = (
                    provider == "azure"
                    or (provider == "auto" and self.config.azure_openai.endpoint)
                )

                if use_azure and AzureOpenAI is not None:
                    self.client = AzureOpenAI(
                        api_key=self.config.azure_openai.api_key,
                        api_version=self.config.azure_openai.api_version,
                        azure_endpoint=self.config.azure_openai.endpoint,
                        timeout=timeout,
                        max_retries=max_retries,
                    )
                    print(f"[OK] Azure OpenAI client initialized (model: {self.config.azure_openai.model})")
                elif OpenAI is not None:
                    self.client = OpenAI(
                        api_key=self.config.azure_openai.api_key,
                        timeout=timeout,
                        max_retries=max_retries,
                    )
                    print(f"[OK] OpenAI client initialized (model: {self.config.azure_openai.model})")
                else:
                    print("[WARN] openai package not installed - AI features disabled")
            except Exception as e:
                print(f"[ERROR] Failed to initialize OpenAI client: {e}")
                print("Proceeding without AI generation - image alt text will use placeholders.")
        else:
            print("[INFO] No API key configured; image alt text will use placeholders.")

    def _get_cNvPr(self, shape):
        """Get the cNvPr (non-visual properties) element for a shape."""
        el = getattr(shape, "_element", None)
        if el is None:
            return None
        tag = el.tag
        if tag.endswith("sp"):
            nv = el.find(qn("p:nvSpPr"))
        elif tag.endswith("pic"):
            nv = el.find(qn("p:nvPicPr"))
        elif tag.endswith("graphicFrame"):
            nv = el.find(qn("p:nvGraphicFramePr"))
        elif tag.endswith("grpSp"):
            nv = el.find(qn("p:nvGrpSpPr"))
        elif tag.endswith("cxnSp"):  # Connection shapes (LINE/connectors)
            nv = el.find(qn("p:nvCxnSpPr"))
        else:
            return None
        return nv.find(qn("p:cNvPr")) if nv is not None else None

    def get_alt_text(self, shape):
        """Get existing alt text from a shape."""
        cNvPr = self._get_cNvPr(shape)
        if cNvPr is not None:
            return cNvPr.get("descr", "")
        return ""

    def set_alt_text(self, shape, alt_text: str):
        """Set alt text for a shape."""
        cNvPr = self._get_cNvPr(shape)
        if cNvPr is not None:
            cNvPr.set("descr", alt_text)
            return True
        return False

    def set_decorative_flag(self, shape, decorative: bool = True):
        """Set or unset the decorative flag for a shape."""
        cNvPr = self._get_cNvPr(shape)
        if cNvPr is None:
            return False

        # Find or create a:extLst
        extLst = cNvPr.find(qn("a:extLst"))
        if extLst is None:
            extLst = OxmlElement("a:extLst")
            cNvPr.append(extLst)

        # Find or create extension for decorative
        ext_el = None
        for ext in extLst.findall(qn("a:ext")):
            if ext.get("uri") == DECORATIVE_EXT_URI:
                ext_el = ext
                break
        if ext_el is None:
            ext_el = OxmlElement("a:ext")
            ext_el.set("uri", DECORATIVE_EXT_URI)
            extLst.append(ext_el)

        # Find or create adec:decorative
        dec_el = None
        for child in ext_el:
            if child.tag == f"{{{ADEC_NS}}}decorative":
                dec_el = child
                break
        if dec_el is None:
            dec_el = OxmlElement("adec:decorative")
            ext_el.append(dec_el)
        dec_el.set("val", "1" if decorative else "0")

        if decorative:
            # clear alt-text description for decorative items
            cNvPr.set("descr", "")

        return True

    def is_decorative_auto_shape(self, shape):
        """Check if a shape is a decorative auto shape.
        
        NOTE: Shapes with images (picture fill) should NOT be marked as decorative.
        """
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            return False
        
        # Check if shape has an image (picture fill) - if so, NOT decorative
        try:
            if hasattr(shape, 'fill') and hasattr(shape.fill, 'type'):
                from pptx.enum.dml import MSO_FILL_TYPE
                if shape.fill.type == MSO_FILL_TYPE.PICTURE:
                    return False  # Has image, not decorative
        except Exception:
            pass
        
        try:
            decorative_auto_shapes = [
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                MSO_AUTO_SHAPE_TYPE.OVAL,
                MSO_AUTO_SHAPE_TYPE.ROUND_1_RECTANGLE,
                MSO_AUTO_SHAPE_TYPE.ROUND_2_DIAG_RECTANGLE,
                MSO_AUTO_SHAPE_TYPE.ROUND_2_SAME_RECTANGLE,
                MSO_AUTO_SHAPE_TYPE.SNIP_1_RECTANGLE,
                MSO_AUTO_SHAPE_TYPE.SNIP_2_DIAG_RECTANGLE,
                MSO_AUTO_SHAPE_TYPE.SNIP_2_SAME_RECTANGLE,
                MSO_AUTO_SHAPE_TYPE.SNIP_ROUND_RECTANGLE,
            ]
            if hasattr(shape, 'auto_shape_type'):
                return shape.auto_shape_type in decorative_auto_shapes
        except Exception:
            pass
        return False

    def is_decorative_container(self, shape):
        """Check if a shape is a decorative container.
        
        NOTE: Shapes with images (picture fill) should NOT be marked as decorative.
        """
        if shape.shape_type not in [MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.AUTO_SHAPE]:
            return False
        
        # Check if shape has an image (picture fill) - if so, NOT decorative
        try:
            if hasattr(shape, 'fill') and hasattr(shape.fill, 'type'):
                from pptx.enum.dml import MSO_FILL_TYPE
                if shape.fill.type == MSO_FILL_TYPE.PICTURE:
                    return False  # Has image, not decorative
        except Exception:
            pass
        
        try:
            if getattr(shape, 'text_frame', None):
                text = shape.text_frame.text.strip()
                if not text:  # Empty text frame
                    return True
        except Exception:
            pass
        return False
    
    def is_connector_shape(self, shape):
        """Check if a shape is a connector/arrow."""
        return shape.shape_type == MSO_SHAPE_TYPE.LINE
    
    def generate_connector_alt_text(self, connector, slide):
        """Generate alt-text for a connector by finding connected shapes and their text.
        
        If the connector connects two text-containing shapes, describe the connection.
        Otherwise, mark it as decorative.
        
        Returns:
            tuple: (alt_text: str or None, is_decorative: bool)
        """
        try:
            # Get connector position
            conn_left = getattr(connector, 'left', 0)
            conn_top = getattr(connector, 'top', 0)
            conn_width = getattr(connector, 'width', 0)
            conn_height = getattr(connector, 'height', 0)
            conn_right = conn_left + conn_width
            conn_bottom = conn_top + conn_height
            
            # Find shapes near the connector's endpoints
            nearby_shapes = []
            threshold = Inches(0.5)  # shapes within 0.5 inches of connector endpoints
            
            for shape in slide.shapes:
                if shape == connector:
                    continue
                if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                    continue
                    
                shape_left = getattr(shape, 'left', 0)
                shape_top = getattr(shape, 'top', 0)
                shape_width = getattr(shape, 'width', 0)
                shape_height = getattr(shape, 'height', 0)
                shape_right = shape_left + shape_width
                shape_bottom = shape_top + shape_height
                
                # Check if shape is near the start or end of the connector
                near_start = (abs(shape_right - conn_left) < threshold or abs(shape_left - conn_left) < threshold) and \
                            (abs(shape_bottom - conn_top) < threshold or abs(shape_top - conn_top) < threshold)
                near_end = (abs(shape_right - conn_right) < threshold or abs(shape_left - conn_right) < threshold) and \
                          (abs(shape_bottom - conn_bottom) < threshold or abs(shape_top - conn_bottom) < threshold)
                
                if near_start or near_end:
                    # Extract text from shape
                    text = ""
                    if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                        text = shape.text_frame.text.strip()
                    
                    if text:
                        position = "start" if near_start else "end"
                        nearby_shapes.append((position, text[:50]))  # Limit to 50 chars
            
            # If we found exactly 2 shapes (one at each end), create descriptive alt-text
            if len(nearby_shapes) == 2:
                start_text = next((text for pos, text in nearby_shapes if pos == "start"), None)
                end_text = next((text for pos, text in nearby_shapes if pos == "end"), None)
                
                if start_text and end_text:
                    return (f"Arrow connecting '{start_text}' to '{end_text}'", False)
                elif start_text or end_text:
                    # Only one endpoint has text
                    text = start_text or end_text
                    return (f"Arrow from/to '{text}'", False)
            
            # Otherwise, mark as decorative
            return (None, True)
            
        except Exception as e:
            # On error, mark as decorative
            return (None, True)

    def _append_order_to_title_attr(self, shape, order_index, prefix="ord"):
        """Append a short order tag into the shape's cNvPr@title attribute.

        This mirrors the behavior in the legacy scripts: it stores a tag like
        "ord 1" in the `title` metadata so reading-order can be inspected later.
        """
        cNvPr = self._get_cNvPr(shape)
        if cNvPr is None:
            return False
        existing = cNvPr.get("title") or ""
        tag = f"{prefix} {order_index}"
        if tag in existing:
            return True
        new_title = (existing + (" | " if existing else "") + tag).strip()
        cNvPr.set("title", new_title)
        return True

    def annotate_reading_order(self, slide):
        """Reorder shapes on a slide to match accessibility reading order.

        Top-to-bottom, left-to-right heuristic is used (sort by `top`, then `left`).
        
        IMPORTANT: This physically reorders the shape XML elements in the slide's <p:spTree>
        to match the reading order. No visible annotations are added to the shapes.
        
        Returns:
            Dict with details: {'shapes_reordered': int, 'shapes': [list of shape names], 'xml_reordered': bool}
        """
        reordered_shapes = []
        xml_reordered = False
        
        # Top-level shapes - sort by position
        sortable = []
        for shp in slide.shapes:
            try:
                t = int(getattr(shp, "top", 0)); l = int(getattr(shp, "left", 0))
                sortable.append((t, l, shp))
            except Exception:
                continue
        sortable.sort(key=lambda x: (x[0], x[1]))
        
        # Track reordered shapes (no visible annotation)
        for idx, (_, __, shp) in enumerate(sortable, start=1):
            try:
                shape_name = getattr(shp, 'name', f'Shape {idx}')
                reordered_shapes.append(f"{shape_name}")
            except Exception:
                pass

        # REORDER THE XML ELEMENTS to match sorted order
        try:
            xml_reordered = self._reorder_shape_xml_elements(slide, sortable)
        except Exception as e:
            print(f"    [WARN] Could not reorder XML elements: {e}")

        # For each group, reorder internal shapes
        group_count = 0
        for shp in slide.shapes:
            try:
                if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
                    group_count += 1
                    inner = []
                    for g in shp.shapes:
                        try:
                            t = int(getattr(g, "top", 0)); l = int(getattr(g, "left", 0))
                            inner.append((t, l, g))
                        except Exception:
                            continue
                    inner.sort(key=lambda x: (x[0], x[1]))
                    for idx, (_, __, g) in enumerate(inner, start=1):
                        try:
                            g_name = getattr(g, 'name', f'Group Item {idx}')
                            reordered_shapes.append(f"  - {g_name}")
                        except Exception:
                            pass
            except Exception:
                continue
        
        return {
            'annotated_count': len(reordered_shapes),  # Keep same key name for compatibility
            'shapes': reordered_shapes,
            'groups_processed': group_count,
            'xml_reordered': xml_reordered
        }
    
    def _reorder_shape_xml_elements(self, slide, sorted_shapes):
        """Physically reorder shape XML elements in the slide's <p:spTree> to match reading order.
        
        Args:
            slide: The slide object
            sorted_shapes: List of (top, left, shape) tuples in desired order
            
        Returns:
            bool: True if reordering succeeded, False otherwise
        """
        try:
            # Get the slide's XML element
            slide_elem = slide.element
            sp_tree = slide_elem.cSld.spTree
            
            # Find non-visual properties element (nvGrpSpPr) - this should stay first
            nvGrpSpPr_elem = None
            for child in sp_tree:
                if child.tag.endswith('}nvGrpSpPr'):
                    nvGrpSpPr_elem = child
                    break
            
            # Find grpSpPr element - this should stay second  
            grpSpPr_elem = None
            for child in sp_tree:
                if child.tag.endswith('}grpSpPr'):
                    grpSpPr_elem = child
                    break
            
            # Build set of shape XML elements for lookup
            shape_elements = set()
            xml_elements_in_order = []
            for shape_tuple in sorted_shapes:
                shape = shape_tuple[2]
                # Get the shape's XML element
                elem = getattr(shape, 'element', None)
                if elem is None:
                    elem = getattr(shape, '_element', None)
                if elem is not None:
                    shape_elements.add(elem)
                    xml_elements_in_order.append(elem)

            # Collect non-shape elements (e.g. mc:AlternateContent, extLst)
            # that python-pptx doesn't expose as shapes but PowerPoint requires
            non_shape_elements = []
            for child in list(sp_tree):
                if child not in [nvGrpSpPr_elem, grpSpPr_elem] and child not in shape_elements:
                    non_shape_elements.append(child)

            # Remove all shape elements from tree (except nvGrpSpPr and grpSpPr)
            elements_to_remove = []
            for child in list(sp_tree):
                if child not in [nvGrpSpPr_elem, grpSpPr_elem]:
                    elements_to_remove.append(child)

            for elem in elements_to_remove:
                sp_tree.remove(elem)

            # Re-add shape elements in sorted order
            for xml_elem in xml_elements_in_order:
                sp_tree.append(xml_elem)

            # Re-add non-shape elements (AlternateContent, etc.) at the end
            for xml_elem in non_shape_elements:
                sp_tree.append(xml_elem)
            
            return True
            
        except Exception as e:
            # If reordering fails, shapes will still have "ord N" annotations
            print(f"    XML reordering exception: {e}")
            import traceback
            traceback.print_exc()
            return False

    def extract_image_data(self, shape):
        """Extract image data from a shape for AI processing using multiple methods."""
        import traceback
        
        shape_name = getattr(shape, 'name', 'Unknown')
        shape_type = getattr(shape, 'shape_type', 'Unknown')
        
        # Method 1: Direct image access (works for PICTURE shapes)
        try:
            if hasattr(shape, 'image') and shape.image:
                return base64.b64encode(shape.image.blob).decode('utf-8')
        except Exception as e:
            print(f"[ERROR] Method 1 (direct image) failed for shape '{shape_name}' type {shape_type}: {e}")
            print(f"[TRACE] {traceback.format_exc()}")
        
        # Method 2: Picture fill extraction (for AUTO_SHAPE with picture fill)
        try:
            if (hasattr(shape, 'fill') and hasattr(shape.fill, 'type')):
                from pptx.enum.dml import MSO_FILL_TYPE
                if shape.fill.type == MSO_FILL_TYPE.PICTURE:
                    # Access the shape's XML element to find the blip embed relationship
                    if hasattr(shape, 'element'):
                        shape_element = shape.element
                        
                        # Recursively search for blip elements with embed attribute
                        def find_blip_embed_id(elem):
                            # Check if this is a blip element with embed attribute
                            if 'blip' in elem.tag and elem.tag.endswith('blip'):
                                for attr_name, attr_value in elem.attrib.items():
                                    if 'embed' in attr_name:
                                        return attr_value
                            
                            # Search children recursively
                            for child in elem:
                                result = find_blip_embed_id(child)
                                if result:
                                    return result
                            return None
                        
                        embed_id = find_blip_embed_id(shape_element)
                        if embed_id:
                            # Get the image part using the relationship
                            part = shape.part
                            if hasattr(part, 'rels') and embed_id in part.rels:
                                rel = part.rels[embed_id]
                                if hasattr(rel, 'target_part'):
                                    image_part = rel.target_part
                                    if hasattr(image_part, 'blob'):
                                        return base64.b64encode(image_part.blob).decode('utf-8')
        except Exception as e:
            print(f"Method 2 (picture fill) failed: {e}")
        
        # Method 3: XML parsing approach (iterative search for blip elements)
        try:
            el = getattr(shape, "_element", None)
            if el is not None:
                for child in el.iter():
                    if child.tag.endswith('blip'):
                        for attr_name, attr_val in child.attrib.items():
                            if attr_name.endswith('embed'):
                                part = shape.part
                                if hasattr(part, 'rels') and attr_val in part.rels:
                                    rel = part.rels[attr_val]
                                    if hasattr(rel, 'target_part') and hasattr(rel.target_part, 'blob'):
                                        return base64.b64encode(rel.target_part.blob).decode('utf-8')
        except Exception as e:
            print(f"Method 3 (XML parsing) failed: {e}")
        
        # Method 4: Original XPath approach with related_parts
        try:
            if hasattr(shape, '_pic') and hasattr(shape._pic, 'xpath'):
                blip_els = shape._pic.xpath('.//a:blip[@r:embed]')
                if blip_els:
                    embed_id = blip_els[0].get(qn('r:embed'))
                    if embed_id:
                        part = shape.part
                        if hasattr(part, 'related_parts') and embed_id in part.related_parts:
                            image_part = part.related_parts[embed_id]
                            if hasattr(image_part, 'blob'):
                                return base64.b64encode(image_part.blob).decode('utf-8')
        except Exception as e:
            print(f"Method 4 (XPath related_parts) failed: {e}")
        
        # If all methods fail, return None
        print(f"Warning: Could not extract image data from shape using any method")
        return None

    def _resize_image_for_api(self, image_data_b64: str, max_size: int = 2048) -> str:
        """Resize image to reduce memory usage while maintaining quality.
        
        Args:
            image_data_b64: Base64 encoded image data
            max_size: Maximum width or height in pixels (default 2048)
            
        Returns:
            Base64 encoded resized image, or original if resize fails
        """
        if not Image:
            return image_data_b64  # PIL not available, return original
            
        try:
            # Decode base64 to bytes
            image_bytes = base64.b64decode(image_data_b64)
            img = Image.open(io.BytesIO(image_bytes))
            
            # Skip if already small enough
            if img.width <= max_size and img.height <= max_size:
                return image_data_b64
                
            # Calculate new dimensions maintaining aspect ratio
            if img.width > img.height:
                new_width = max_size
                new_height = int((max_size / img.width) * img.height)
            else:
                new_height = max_size
                new_width = int((max_size / img.height) * img.width)
                
            # Resize using high-quality resampling
            img_resized = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
            
            # Convert back to base64
            buffer = io.BytesIO()
            img_format = img.format if img.format else 'PNG'
            img_resized.save(buffer, format=img_format, optimize=True, quality=85)
            resized_b64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
            
            # Log size reduction
            original_size = len(image_data_b64)
            new_size = len(resized_b64)
            if original_size > new_size:
                reduction = ((original_size - new_size) / original_size) * 100
                print(f"    Resized image: {img.width}x{img.height} -> {new_width}x{new_height} ({reduction:.1f}% smaller)")
            
            return resized_b64
        except Exception as e:
            print(f"    [WARN] Image resize failed: {e}, using original")
            return image_data_b64

    @staticmethod
    def _strip_alt_text_prefix(text: str) -> str:
        """Strip 'Alt text:', 'Alt-text:', etc. prefixes that AI models sometimes add."""
        import re
        stripped = re.sub(r'^(?:alt[\s\-]*text)\s*[:;.\-]\s*', '', text, count=1, flags=re.IGNORECASE)
        return stripped

    def generate_alt_text_with_ai(self, image_data: str, context: str = "") -> str:
        """Generate alt text using Azure OpenAI with caching, retry, and memory optimization."""
        import time as _time

        if not self.client:
            return "Image content - AI description not available"

        # Resize image to reduce memory usage
        image_data = self._resize_image_for_api(image_data, max_size=2048)

        # Check cache if enabled
        enable_caching = getattr(self.config.processing, 'enable_image_caching', True)
        if enable_caching:
            # Use first 100 chars of base64 as cache key (enough to identify duplicates)
            cache_key = image_data[:100] if len(image_data) > 100 else image_data
            if cache_key in self._image_cache:
                return self._image_cache[cache_key]

        prompt = self.config.azure_openai.prompt_template.format(
            max_length=self.config.processing.max_alt_text_length
        )
        if context:
            prompt = f"{prompt} Context: {context}"

        # Retry with exponential backoff for rate limits
        max_attempts = 3
        for attempt in range(max_attempts):
            try:
                response = self.client.chat.completions.create(
                    model=self.config.azure_openai.model,
                    messages=[
                        {
                            "role": "user",
                            "content": [
                                {"type": "text", "text": prompt},
                                {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{image_data}"}}
                            ]
                        }
                    ],
                    max_tokens=self.config.azure_openai.max_tokens,
                    temperature=self.config.azure_openai.temperature,
                )
                alt_text = response.choices[0].message.content.strip()
                alt_text = self._strip_alt_text_prefix(alt_text)

                # Truncate if too long
                if len(alt_text) > self.config.processing.max_alt_text_length:
                    alt_text = alt_text[:self.config.processing.max_alt_text_length - 3] + "..."

                # Cache the result
                if enable_caching:
                    cache_key = image_data[:100] if len(image_data) > 100 else image_data
                    self._image_cache[cache_key] = alt_text

                return alt_text
            except Exception as e:
                error_msg = str(e)

                # Check for Azure content filter violations (common with medical images)
                if 'content_filter' in error_msg.lower() or 'responsibleaipolicyviolation' in error_msg.lower():
                    print(f"    [WARN] Content filter triggered (likely medical/surgical content): {error_msg[:150]}")
                    return "Medical/technical image content - filtered by content policy"

                # Check for rate limiting - retry with backoff
                if 'rate' in error_msg.lower() or '429' in error_msg:
                    if attempt < max_attempts - 1:
                        wait_time = (attempt + 1) * 5  # 5s, 10s
                        print(f"    [WARN] Rate limited, waiting {wait_time}s before retry (attempt {attempt + 1}/{max_attempts})")
                        _time.sleep(wait_time)
                        continue
                    print(f"    [WARN] Rate limited after {max_attempts} attempts: {error_msg[:150]}")
                    return "Image content - rate limit exceeded, try again later"

                # Check for timeout - retry once
                if 'timeout' in error_msg.lower() or 'timed out' in error_msg.lower():
                    if attempt < max_attempts - 1:
                        print(f"    [WARN] API timeout, retrying (attempt {attempt + 1}/{max_attempts})")
                        continue
                    print(f"    [WARN] API timeout after {max_attempts} attempts: {error_msg[:150]}")
                    return "Image content - API timeout"

                # Generic error - don't retry
                print(f"    [ERROR] AI alt-text generation failed: {error_msg[:200]}")
                return "Image content - AI description failed"

        return "Image content - AI description failed"

    def generate_alt_text_for_text_shape(self, text_content: str, slide_num: int) -> str:
        """Generate alt text for a shape containing text using AI."""
        if not self.client:
            # Fallback to simple description if AI not available
            preview = text_content[:50] + "..." if len(text_content) > 50 else text_content
            return f"Text box: {preview}"

        try:
            prompt = (
                f"Describe this text element from a PowerPoint slide concisely "
                f"in {self.config.processing.max_alt_text_length} characters or less. "
                f"The text content is: '{text_content}'. "
                f"Describe what this text element represents in the presentation context. "
                f"Do not start with 'Alt text' or 'Alt-text'."
            )

            response = self.client.chat.completions.create(
                model=self.config.azure_openai.model,
                messages=[
                    {
                        "role": "user",
                        "content": prompt
                    }
                ],
                max_tokens=self.config.azure_openai.max_tokens,
                temperature=self.config.azure_openai.temperature,
            )
            alt_text = response.choices[0].message.content.strip()
            alt_text = self._strip_alt_text_prefix(alt_text)

            # Truncate if too long
            if len(alt_text) > self.config.processing.max_alt_text_length:
                alt_text = alt_text[:self.config.processing.max_alt_text_length - 3] + "..."

            return alt_text
        except Exception as e:
            print(f"Warning: AI alt-text generation for text shape failed: {e}")
            # Fallback to simple description
            preview = text_content[:50] + "..." if len(text_content) > 50 else text_content
            return f"Text box: {preview}"

    def process_shape(self, shape, slide_num: int, shape_idx: int, slide=None) -> Dict[str, Any]:
        """Process a single shape for alt-text."""
        import traceback
        
        shape_name = getattr(shape, 'name', 'Unknown')
        result = {
            'shape_idx': shape_idx,
            'shape_name': shape_name,
            'shape_type': str(shape.shape_type),
            'processed': False,
            'alt_text': '',
            'decorative': False,
            'error': None
        }

        try:
            print(f"[DEBUG] Processing slide {slide_num}, shape {shape_idx} '{shape_name}' type {shape.shape_type}")
            existing_alt = self.get_alt_text(shape)
            
            # Skip if alt text exists and not forcing regeneration
            if existing_alt and not self.config.processing.force_regenerate:
                result['alt_text'] = existing_alt
                result['processed'] = True
                return result

            # Handle connectors/arrows
            if self.is_connector_shape(shape):
                if slide:
                    alt_text, is_decorative = self.generate_connector_alt_text(shape, slide)
                    if is_decorative:
                        self.set_decorative_flag(shape, True)
                        result['decorative'] = True
                    elif alt_text:
                        self.set_alt_text(shape, alt_text)
                        result['alt_text'] = alt_text
                else:
                    # No slide context, mark as decorative
                    self.set_decorative_flag(shape, True)
                    result['decorative'] = True
                result['processed'] = True
                return result

            # Check if it's decorative
            if (self.is_decorative_auto_shape(shape) or 
                self.is_decorative_container(shape)):
                self.set_decorative_flag(shape, True)
                result['decorative'] = True
                result['processed'] = True
                return result

            # If configured, skip adding alt-text to textboxes/auto-shapes
            # (this keeps images, tables, charts, etc. processed as before)
            if getattr(self.config.processing, 'images_and_decorators_only', False):
                if shape.shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.TEXT_BOX]:
                    return result

            # Process images
            if (shape.shape_type == MSO_SHAPE_TYPE.PICTURE and 
                self.config.processing.process_images):
                image_data = self.extract_image_data(shape)
                if image_data:
                    context = f"This is from slide {slide_num}"
                    alt_text = self.generate_alt_text_with_ai(image_data, context)
                    self.set_alt_text(shape, alt_text)
                    result['alt_text'] = alt_text
                    result['processed'] = True

            # Process other shapes
            elif self.config.processing.process_shapes:
                if shape.shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.TEXT_BOX]:
                    # Skip text boxes if configured (text is already accessible to screen readers)
                    if getattr(self.config.processing, 'skip_text_boxes', False) and shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                        result['processed'] = False
                        return result
                    
                    # Check if shape has text content
                    text_content = ""
                    if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                        text_content = shape.text_frame.text.strip()

                    # If configured to skip objects that contain text, and this shape has text, skip
                    if getattr(self.config.processing, 'skip_objects_with_text', False) and text_content:
                        result['processed'] = False
                        return result
                    
                    # Check if shape also has an image (filled with picture)
                    image_data = None
                    try:
                        if hasattr(shape, 'fill') and hasattr(shape.fill, 'type'):
                            from pptx.enum.dml import MSO_FILL_TYPE
                            if shape.fill.type == MSO_FILL_TYPE.PICTURE:
                                image_data = self.extract_image_data(shape)
                    except:
                        pass
                    
                    # Generate alt text based on available content
                    if image_data and text_content:
                        # Shape has both image and text - use AI with both as context
                        context = f"Shape on slide {slide_num} containing text: '{text_content}'"
                        alt_text = self.generate_alt_text_with_ai(image_data, context)
                    elif image_data:
                        # Shape has only image
                        context = f"Shape on slide {slide_num}"
                        alt_text = self.generate_alt_text_with_ai(image_data, context)
                    elif text_content:
                        # Shape has only text - use AI to generate descriptive alt text
                        alt_text = self.generate_alt_text_for_text_shape(text_content, slide_num)
                    else:
                        # Shape has neither - use basic description
                        alt_text = f"Shape on slide {slide_num}"
                    
                    self.set_alt_text(shape, alt_text)
                    result['alt_text'] = alt_text
                    result['processed'] = True

        except Exception as e:
            result['error'] = str(e)
            print(f"Error processing shape {shape_idx} on slide {slide_num}: {e}")

        return result

    def generate_slide_title(self, slide, slide_num: int) -> str:
        """Generate a title for a slide based on its content."""
        try:
            text_shapes = []
            image_count = 0
            table_count = 0
            
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                    text_shapes.append(shape.text_frame.text.strip())
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_count += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_count += 1

            if text_shapes:
                # Find the longest text that might be a title
                longest = max(text_shapes, key=len)
                if len(longest) > 50:
                    if '.' in longest[:50]:
                        longest = longest[:longest.find('.', 30)]
                    elif ',' in longest[:50]:
                        longest = longest[:longest.find(',', 30)]
                    else:
                        longest = longest[:47] + "..."
                heuristic_title = longest or f"Slide {slide_num}"
                # If AI client is available, prefer AI-generated concise title
                if self.client:
                    try:
                        prompt = (
                            f"Generate a short slide title (max 50 chars) for slide {slide_num} based on this text: '{heuristic_title}'"
                        )
                        response = self.client.chat.completions.create(
                            model=self.config.azure_openai.model,
                            messages=[{"role": "user", "content": prompt}],
                            max_tokens=20,
                            temperature=0.0,
                        )
                        ai_title = (response.choices[0].message.content or "").strip()
                        if ai_title:
                            return ai_title
                    except Exception:
                        pass
                return heuristic_title
            
            if image_count and table_count:
                return f"Slide {slide_num}: Visual Content and Data"
            if image_count:
                return f"Slide {slide_num}: Visual Content"
            if table_count:
                return f"Slide {slide_num}: Data and Tables"
            return f"Slide {slide_num}"
        except Exception:
            return f"Slide {slide_num}"

    def get_existing_slide_title(self, slide) -> str:
        """Get the existing title from a slide if it exists."""
        # Check title placeholders first
        for shape in slide.shapes:
            if getattr(shape, 'is_placeholder', False):
                try:
                    if shape.placeholder_format and shape.placeholder_format.type.name in ['TITLE', 'CENTER_TITLE']:
                        if getattr(shape, 'text_frame', None):
                            existing_text = shape.text_frame.text.strip()
                            if existing_text:
                                return existing_text
                except Exception:
                    pass
        
        # Check top-most text box as potential title
        top_candidate = None
        top_y = None
        for shape in slide.shapes:
            if getattr(shape, 'text_frame', None):
                y = getattr(shape, 'top', Inches(10))
                if top_y is None or y < top_y:
                    top_y = y
                    top_candidate = shape
        
        if top_candidate is not None:
            try:
                existing_text = top_candidate.text_frame.text.strip()
                if existing_text:
                    return existing_text
            except Exception:
                pass
        
        return ""

    def set_slide_title(self, slide, title_text: str) -> bool:
        """Set the title of a slide."""
        # 1) Fill an existing title placeholder
        for shape in slide.shapes:
            if getattr(shape, 'is_placeholder', False):
                try:
                    if shape.placeholder_format and shape.placeholder_format.type.name in ['TITLE', 'CENTER_TITLE']:
                        if getattr(shape, 'text_frame', None):
                            shape.text_frame.clear()
                            p = shape.text_frame.paragraphs[0]
                            run = p.add_run()
                            run.text = title_text
                            run.font.bold = True
                            if not run.font.size:
                                run.font.size = Pt(24)
                            return True
                except Exception:
                    pass
        
        # 2) Heuristic overwrite of top-area textbox
        top_candidate = None
        top_y = None
        for shape in slide.shapes:
            if getattr(shape, 'text_frame', None):
                y = getattr(shape, 'top', Inches(10))
                if top_y is None or y < top_y:
                    top_y = y
                    top_candidate = shape
        if top_candidate is not None:
            try:
                tf = top_candidate.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = title_text
                run.font.bold = True
                if not run.font.size:
                    run.font.size = Pt(24)
                return True
            except Exception:
                pass
        
        # 3) Add a new textbox at the top
        try:
            left = Inches(0.5)
            top = Inches(0.25)
            width = Inches(9)
            height = Inches(0.7)
            tb = slide.shapes.add_textbox(left, top, width, height)
            tf = tb.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title_text
            run.font.bold = True
            run.font.size = Pt(24)
            return True
        except Exception:
            pass
        
        return False

    def process_presentation(self, input_path: str, output_path: str = None, progress_callback=None) -> Dict[str, Any]:
        """Process a PowerPoint presentation to add alt-text and slide titles.

        Args:
            progress_callback: Optional callable(current_slide, total_slides, shapes_processed, elapsed, remaining)
        """
        import gc
        import time as _time

        # Refresh configuration at runtime to pick up external test overrides
        try:
            self.config = get_config()
        except Exception:
            pass
        # Reset seen titles per presentation to avoid cross-request bleed
        self._seen_titles = {}
        # Reset image cache per presentation
        self._image_cache = {}
        if output_path is None:
            base = Path(input_path).stem
            output_path = f"{base}_with_alttext_v2.pptx"

        # Backup original if configured
        if self.config.processing.backup_originals:
            backup_path = f"{input_path}.backup"
            if not os.path.exists(backup_path):
                shutil.copy2(input_path, backup_path)

        try:
            prs = Presentation(input_path)
            total_slides = len(prs.slides)
            results = {
                'input_file': input_path,
                'output_file': output_path,
                'total_slides': total_slides,
                'processed_slides': 0,
                'processed_shapes': 0,
                'errors': [],
                'slide_details': []
            }

            # Adaptive concurrency: reduce thread count for large presentations
            # to avoid Azure OpenAI rate limiting and memory pressure
            max_workers = self.config.processing.max_concurrent_api_calls
            is_large_file = total_slides > 50
            if is_large_file:
                max_workers = min(max_workers, 3)
                print(f"[INFO] Large presentation ({total_slides} slides) - reducing concurrency to {max_workers} threads")

            processing_start = _time.time()

            # Process each slide
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_start = _time.time()
                slide_result = {
                    'slide_num': slide_num,
                    'shapes_processed': 0,
                    'title_set': False,
                    'shapes': []
                }
                # Determine existing title once for this slide
                existing_title = self.get_existing_slide_title(slide)

                # Title handling logic: only set a title if missing or forced; only deduplicate when repeated
                if self.config.processing.process_slide_titles:
                    if not existing_title or self.config.processing.force_regenerate_slide_titles:
                        title = self.generate_slide_title(slide, slide_num)

                        if self.config.processing.deduplicate_titles:
                            if title in self._seen_titles:
                                self._seen_titles[title] += 1
                                title = f"{title} - Slide {slide_num}"
                            else:
                                self._seen_titles[title] = 1

                        if self.set_slide_title(slide, title):
                            slide_result['title_set'] = True
                            slide_result['title'] = title
                    else:
                        slide_result['title'] = existing_title
                        if self.config.processing.deduplicate_titles:
                            if existing_title in self._seen_titles:
                                self._seen_titles[existing_title] += 1
                                new_title = f"{existing_title} - Slide {slide_num}"
                                if self.set_slide_title(slide, new_title):
                                    slide_result['title_set'] = True
                                    slide_result['title'] = new_title
                            else:
                                self._seen_titles[existing_title] = 1
                else:
                    # If title processing disabled, still track seen titles to avoid cross-run collisions
                    if existing_title and self.config.processing.deduplicate_titles:
                        self._seen_titles[existing_title] = self._seen_titles.get(existing_title, 0) + 1

                # Process shapes on the slide
                shapes_to_process = [(i, shape) for i, shape in enumerate(slide.shapes)]

                if self.config.processing.enable_multithreading and len(shapes_to_process) > 1:
                    # Use threading for shape processing
                    with ThreadPoolExecutor(max_workers=max_workers) as executor:
                        future_to_shape = {
                            executor.submit(self.process_shape, shape, slide_num, i, slide): (i, shape)
                            for i, shape in shapes_to_process
                        }

                        for future in as_completed(future_to_shape):
                            try:
                                shape_result = future.result()
                                slide_result['shapes'].append(shape_result)
                                if shape_result['processed']:
                                    slide_result['shapes_processed'] += 1
                            except Exception as e:
                                results['errors'].append(f"Slide {slide_num}: {e}")
                else:
                    # Process shapes sequentially
                    for i, shape in shapes_to_process:
                        shape_result = self.process_shape(shape, slide_num, i, slide)
                        slide_result['shapes'].append(shape_result)
                        if shape_result['processed']:
                            slide_result['shapes_processed'] += 1

                # After processing all shapes on this slide, annotate reading order
                reading_order_result = self.annotate_reading_order(slide)
                slide_result['reading_order'] = reading_order_result

                # Log reading order changes
                if reading_order_result['annotated_count'] > 0:
                    print(f"  Reading order applied to {reading_order_result['annotated_count']} shapes")
                    for shape_info in reading_order_result['shapes'][:5]:  # Show first 5
                        print(f"    {shape_info}")
                    if len(reading_order_result['shapes']) > 5:
                        print(f"    ... and {len(reading_order_result['shapes']) - 5} more")

                results['slide_details'].append(slide_result)
                results['processed_slides'] += 1
                results['processed_shapes'] += slide_result['shapes_processed']

                # Progress reporting
                total_elapsed = _time.time() - processing_start
                avg_per_slide = total_elapsed / slide_num
                remaining = avg_per_slide * (total_slides - slide_num)

                # Call progress callback every slide
                if progress_callback:
                    try:
                        progress_callback(slide_num, total_slides, results['processed_shapes'], total_elapsed, remaining)
                    except Exception:
                        pass

                # Console progress logging for large files
                slide_elapsed = _time.time() - slide_start
                if is_large_file and (slide_num % 10 == 0 or slide_num == total_slides):
                    print(f"[PROGRESS] Slide {slide_num}/{total_slides} "
                          f"({slide_result['shapes_processed']} shapes, {slide_elapsed:.1f}s) "
                          f"- elapsed {total_elapsed:.0f}s, ~{remaining:.0f}s remaining")

                # Periodic memory cleanup for large files
                if is_large_file and slide_num % 25 == 0:
                    gc.collect()

            # Save the presentation
            print(f"[INFO] Saving presentation ({total_slides} slides)...")
            prs.save(output_path)
            results['success'] = True

            total_time = _time.time() - processing_start
            print(f"[OK] Processed {results['processed_slides']} slides with {results['processed_shapes']} shapes in {total_time:.1f}s")
            print(f"[OK] Saved to: {output_path}")

            return results

        except Exception as e:
            error_msg = f"Failed to process presentation: {e}"
            print(f"[ERROR] {error_msg}")
            return {
                'success': False,
                'error': error_msg,
                'input_file': input_path
            }
