"""
Accessibility scoring utility for PowerPoint presentations (V2)
Based on accessibility_scorer_no_com.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from typing import Dict, Any, List

from ..config import get_config


def _get_cNvPr(shape):
    """
    Return the <p:cNvPr> element for a shape regardless of its concrete type.
    Falls back to None when not found.
    """
    el = getattr(shape, "_element", None)
    if el is None:
        return None
    
    tag = el.tag
    if tag.endswith("sp"):
        nv = el.find(qn("p:nvSpPr"))
    elif tag.endswith("pic"):
        nv = el.find(qn("p:nvPicPr"))
    elif tag.endswith("graphicFrame"):
        nv = el.find(qn("p:nvGraphicFramePr"))
    elif tag.endswith("grpSp"):
        nv = el.find(qn("p:nvGrpSpPr"))
    else:
        nv = None
    
    if nv is None:
        return None
    return nv.find(qn("p:cNvPr"))


def get_alt_text(shape):
    """
    Retrieve alternative text (description) for any shape.
    Tries python-pptx API if present, then falls back to reading cNvPr@descr.
    """
    # Some python-pptx versions expose shape.alternative_text; many do not.
    alt = getattr(shape, "alternative_text", None)
    if isinstance(alt, str) and alt.strip():
        return alt.strip()
    
    cNvPr = _get_cNvPr(shape)
    if cNvPr is None:
        return ""
    return (cNvPr.get("descr") or "").strip()


class AccessibilityScorer:
    """Calculate accessibility scores for PowerPoint presentations."""

    def __init__(self):
        self.config = get_config()

    def calculate_accessibility_score(self, pptx_path: str) -> Dict[str, Any]:
        """
        Calculate accessibility score based on alt-text coverage.
        Returns a dictionary with detailed metrics and overall score.
        """
        pptx_path = Path(pptx_path)

        if not pptx_path.exists():
            return {"error": f"File not found: {pptx_path}"}

        try:
            prs = Presentation(pptx_path)

            # Initialize counters
            total_images = 0
            images_with_alttext = 0
            total_shapes = 0
            shapes_with_alttext = 0
            total_groups = 0
            slide_count = len(prs.slides)

            slide_details = []

            # Process each slide
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_images = 0
                slide_images_with_alt = 0
                slide_shapes = 0
                slide_shapes_with_alt = 0
                slide_groups = 0

                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        slide_groups += 1
                        total_groups += 1

                        for sub_shape in shape.shapes:
                            if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                slide_images += 1
                                total_images += 1
                                if get_alt_text(sub_shape):
                                    slide_images_with_alt += 1
                                    images_with_alttext += 1
                            else:
                                slide_shapes += 1
                                total_shapes += 1
                                if get_alt_text(sub_shape):
                                    slide_shapes_with_alt += 1
                                    shapes_with_alttext += 1

                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        slide_images += 1
                        total_images += 1
                        if get_alt_text(shape):
                            slide_images_with_alt += 1
                            images_with_alttext += 1

                    elif shape.shape_type in [
                        MSO_SHAPE_TYPE.AUTO_SHAPE,
                        MSO_SHAPE_TYPE.TEXT_BOX,
                        MSO_SHAPE_TYPE.CHART,
                        MSO_SHAPE_TYPE.TABLE,
                        MSO_SHAPE_TYPE.SMART_ART,
                        MSO_SHAPE_TYPE.MEDIA,
                    ]:
                        slide_shapes += 1
                        total_shapes += 1
                        if get_alt_text(shape):
                            slide_shapes_with_alt += 1
                            shapes_with_alttext += 1

                slide_details.append({
                    "slide_number": slide_num,
                    "images": slide_images,
                    "images_with_alt": slide_images_with_alt,
                    "shapes": slide_shapes,
                    "shapes_with_alt": slide_shapes_with_alt,
                    "groups": slide_groups,
                })

            # Calculate percentages
            image_coverage = (images_with_alttext / total_images * 100) if total_images > 0 else 100
            shape_coverage = (shapes_with_alttext / total_shapes * 100) if total_shapes > 0 else 100
            
            # Overall score (weighted average: images are more critical)
            if total_images > 0 and total_shapes > 0:
                overall_score = (image_coverage * 0.7 + shape_coverage * 0.3)
            elif total_images > 0:
                overall_score = image_coverage
            elif total_shapes > 0:
                overall_score = shape_coverage
            else:
                overall_score = 100  # No content to check

            # Determine accessibility level
            if overall_score >= 90:
                accessibility_level = "Excellent"
            elif overall_score >= 80:
                accessibility_level = "Good"
            elif overall_score >= 60:
                accessibility_level = "Fair"
            elif overall_score >= 40:
                accessibility_level = "Poor"
            else:
                accessibility_level = "Very Poor"

            # Check if meets target threshold
            meets_target = overall_score >= self.config.accessibility.target_score_threshold

            return {
                "file_path": str(pptx_path),
                "slide_count": slide_count,
                "total_images": total_images,
                "images_with_alttext": images_with_alttext,
                "total_shapes": total_shapes,
                "shapes_with_alttext": shapes_with_alttext,
                "total_groups": total_groups,
                "image_coverage_percentage": round(image_coverage, 2),
                "shape_coverage_percentage": round(shape_coverage, 2),
                "overall_accessibility_score": round(overall_score, 2),
                "accessibility_level": accessibility_level,
                "meets_target_threshold": meets_target,
                "target_threshold": self.config.accessibility.target_score_threshold,
                "slide_details": slide_details,
                "summary": {
                    "total_items": total_images + total_shapes,
                    "items_with_alt": images_with_alttext + shapes_with_alttext,
                    "missing_alt_text": (total_images + total_shapes) - (images_with_alttext + shapes_with_alttext)
                }
            }

        except Exception as e:
            return {
                "error": f"Failed to analyze presentation: {str(e)}",
                "file_path": str(pptx_path)
            }

    def generate_accessibility_report(self, score_data: Dict[str, Any]) -> str:
        """Generate a human-readable accessibility report."""
        if "error" in score_data:
            return f"Error: {score_data['error']}"

        report = []
        report.append("# PowerPoint Accessibility Report")
        report.append(f"**File:** {score_data['file_path']}")
        report.append(f"**Analysis Date:** {Path(__file__).stat().st_mtime}")
        report.append("")

        # Overall Summary
        report.append("## Overall Summary")
        report.append(f"- **Accessibility Score:** {score_data['overall_accessibility_score']}% ({score_data['accessibility_level']})")
        report.append(f"- **Target Threshold:** {score_data['target_threshold']}% ({'Met' if score_data['meets_target_threshold'] else 'Not Met'})")
        report.append(f"- **Total Slides:** {score_data['slide_count']}")
        report.append("")

        # Detailed Metrics
        report.append("## Detailed Metrics")
        report.append("### Images")
        report.append(f"- Total Images: {score_data['total_images']}")
        report.append(f"- Images with Alt-Text: {score_data['images_with_alttext']}")
        report.append(f"- Image Coverage: {score_data['image_coverage_percentage']}%")
        report.append("")

        report.append("### Shapes & Objects")
        report.append(f"- Total Shapes: {score_data['total_shapes']}")
        report.append(f"- Shapes with Alt-Text: {score_data['shapes_with_alttext']}")
        report.append(f"- Shape Coverage: {score_data['shape_coverage_percentage']}%")
        report.append("")

        if score_data['total_groups'] > 0:
            report.append(f"### Groups")
            report.append(f"- Total Groups: {score_data['total_groups']}")
            report.append("")

        # Summary Stats
        summary = score_data['summary']
        missing = summary['missing_alt_text']
        if missing > 0:
            report.append("## Issues Found")
            report.append(f"- **{missing} items** are missing alt-text descriptions")
            report.append("- Consider adding descriptive alt-text to improve accessibility")
            report.append("")

        # Slide-by-slide breakdown
        report.append("## Slide-by-Slide Breakdown")
        for slide in score_data['slide_details']:
            slide_num = slide['slide_number']
            slide_issues = []
            
            if slide['images'] > slide['images_with_alt']:
                missing_images = slide['images'] - slide['images_with_alt']
                slide_issues.append(f"{missing_images} image(s) missing alt-text")
            
            if slide['shapes'] > slide['shapes_with_alt']:
                missing_shapes = slide['shapes'] - slide['shapes_with_alt']
                slide_issues.append(f"{missing_shapes} shape(s) missing alt-text")
            
            status = "Issues found" if slide_issues else "All content has alt-text"
            report.append(f"**Slide {slide_num}:** {status}")
            
            if slide_issues:
                for issue in slide_issues:
                    report.append(f"  - {issue}")
            
            report.append("")

        # Recommendations
        if missing > 0:
            report.append("## Recommendations")
            report.append("1. Add descriptive alt-text to all images")
            report.append("2. Add alt-text to complex shapes and charts")
            report.append("3. Mark purely decorative elements as decorative")
            report.append("4. Keep alt-text concise but descriptive (under 180 characters)")
            report.append("5. Re-run this analysis after making changes")

        return "\n".join(report)